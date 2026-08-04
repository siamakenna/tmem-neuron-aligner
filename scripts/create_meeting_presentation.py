from __future__ import annotations

from pathlib import Path

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "outputs"
OUT_DIR.mkdir(exist_ok=True)
FINAL = OUT_DIR / "tmem106b_neuron_alignment_meeting_deck.pptx"

ALL = ROOT / "reports/260213_all_wells_20260623_days8_12_16"
PILOT = ROOT / "reports/260213_pilot_20260623_125859"
SINGLE = PILOT / "single_neuron_examples"

W = 7.5
H = 10.0

DARK = RGBColor(70, 70, 70)
BLACK = RGBColor(20, 20, 20)
MUTED = RGBColor(95, 95, 95)
GREEN = RGBColor(214, 236, 221)
GREEN2 = RGBColor(232, 245, 236)
BLUE = RGBColor(121, 195, 226)
LINE = RGBColor(180, 180, 180)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    summary = pd.read_csv(ALL / "all_wells_summary_stats.csv")
    measurements = pd.read_csv(ALL / "all_wells_mcherry_measurements.csv")
    qc = pd.read_csv(ALL / "all_wells_registration_qc.csv")

    title_slide(prs)
    pipeline_slide(prs)
    all_well_scope_slide(prs, qc)
    mcherry_result_slide(prs, measurements)
    plate_variability_slide(prs)
    single_neuron_slide(prs)
    omezarr_slide(prs)
    next_steps_slide(prs)

    prs.save(FINAL)
    write_speaker_notes()
    print(FINAL)


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def style_slide(slide, page: int, section: str = "TMEM106B neuron alignment"):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    # Pale organic-looking bands, echoing the reference PDF.
    add_shape(slide, MSO_SHAPE.OVAL, -1.15, 0.55, 2.05, 8.9, GREEN, None, 0)
    add_shape(slide, MSO_SHAPE.OVAL, 6.45, 5.85, 1.45, 3.5, GREEN, None, 0)
    add_shape(slide, MSO_SHAPE.OVAL, 6.15, 8.1, 1.1, 1.1, RGBColor(238, 248, 241), None, 0)
    # Light blue dots/circuit accents.
    for x, y, r in [(6.65, 8.95, 0.055), (6.9, 9.22, 0.035), (6.53, 9.35, 0.035)]:
        add_shape(slide, MSO_SHAPE.OVAL, x, y, r, r, BLUE, None, 0)
    plus = add_text(slide, "+", 6.75, 8.73, 0.28, 0.18, 12, BLUE, bold=True, align=PP_ALIGN.CENTER)
    plus.rotation = 0
    add_text(slide, section, 5.45, 9.55, 1.45, 0.18, 6.5, MUTED, align=PP_ALIGN.RIGHT)
    add_text(slide, str(page), 0.52, 9.45, 0.25, 0.18, 7, MUTED)


def add_shape(slide, shape, x, y, w, h, fill, line=LINE, line_width=0.6):
    shp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_width)
    return shp


def add_text(slide, text, x, y, w, h, size, color=BLACK, bold=False, italic=False, font="Georgia", align=None):
    text = str(text)
    if not text:
        text = " "
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_bullets(slide, bullets, x, y, w, h, size=13, color=BLACK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Georgia"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 1.1, 1.15, 5.1, 1.15, 26, DARK, bold=True)
    if subtitle:
        add_text(slide, subtitle, 1.12, 2.2, 4.9, 0.45, 11.5, MUTED, font="Georgia")


def add_image(slide, path: Path, x, y, w, h, crop=False):
    path = ensure_png_or_supported(path)
    if not path.exists():
        return None
    with Image.open(path) as im:
        iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if crop:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if img_ratio >= box_ratio:
        width = w
        height = w / img_ratio
        top = y + (h - height) / 2
        left = x
    else:
        height = h
        width = h * img_ratio
        left = x + (w - width) / 2
        top = y
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def ensure_png_or_supported(path: Path) -> Path:
    return path


def title_slide(prs):
    slide = blank(prs)
    style_slide(slide, 1, "")
    add_shape(slide, MSO_SHAPE.OVAL, -1.25, 0.85, 2.1, 8.2, GREEN, None)
    add_text(slide, "TMEM106B neuron alignment", 1.08, 1.85, 5.2, 1.55, 27, DARK, bold=True)
    add_text(slide, "All-well longitudinal pilot and single-neuron candidate examples", 1.1, 3.38, 4.95, 0.62, 14, BLACK)
    add_text(slide, "260213 Feb recopy dataset | Days 8, 12, 16", 1.1, 4.25, 4.8, 0.35, 10.5, MUTED)
    add_text(slide, "Prepared for PI meeting", 1.1, 7.15, 2.4, 0.3, 10, MUTED)
    add_text(slide, "Key takeaway: the pipeline now scales from one pilot pair to all 192 wells and shows stronger punctate-to-diffuse mCherry redistribution in PLD3+TMEM106B+mCherry wells.", 3.35, 6.55, 2.9, 1.25, 11, BLACK)


def pipeline_slide(prs):
    slide = blank(prs)
    style_slide(slide, 2)
    add_title(slide, "From raw ND2 to trajectories", "Register on 488, measure mCherry as phenotype.")
    add_bullets(
        slide,
        [
            "Raw fluorescence ND2 files are loaded read-only.",
            "488 channel estimates X/Y drift, avoiding bias from mCherry redistribution.",
            "Same transform is applied to 561/mCherry, then puncta/diffuse signal is quantified.",
            "Outputs are CSV/QC figures now; OME-Zarr stores support scalable image review.",
        ],
        1.12,
        2.95,
        2.15,
        3.1,
        12,
    )
    add_image(slide, PILOT / "figures/registration_before_after.png", 3.35, 2.75, 3.15, 3.1)
    add_text(slide, "Example registration QC: before/after 488-channel alignment", 3.42, 5.95, 2.95, 0.42, 8.5, MUTED)


def all_well_scope_slide(prs, qc):
    slide = blank(prs)
    style_slide(slide, 3)
    add_title(slide, "Scale-up: all wells, not just E05/F05", "The E05/F05 pair was only the first pilot.")
    stats = [
        ("192/192", "wells processed"),
        ("576", "ND2 files read"),
        ("576/576", "registration QC pass"),
        ("96", "mCherry-valid wells measured"),
    ]
    for i, (num, label) in enumerate(stats):
        x = 1.05 + (i % 2) * 2.45
        y = 2.75 + (i // 2) * 1.25
        add_text(slide, num, x, y, 1.55, 0.42, 22, DARK, bold=True)
        add_text(slide, label, x, y + 0.42, 1.9, 0.32, 9.5, MUTED)
    add_image(slide, ALL / "figures/all_wells_registration_qc_pass_fraction.png", 1.05, 5.35, 5.15, 2.45)
    add_text(slide, "Rows C/D/G/H/K/L are included for registration QC, but excluded from mCherry puncta/diffuse interpretation.", 1.1, 8.15, 5.5, 0.55, 10.5, BLACK)


def mcherry_result_slide(prs, measurements):
    slide = blank(prs)
    style_slide(slide, 4)
    add_title(slide, "mCherry signal shifts more in TMEM106B wells", "Screening metric: diffuse / punctate mCherry ratio.")
    add_image(slide, ALL / "figures/all_wells_mcherry_condition_summary.png", 0.95, 2.45, 5.85, 3.35)
    add_bullets(
        slide,
        [
            "Primary wells: 3.26 -> 4.53 from Day 8 to Day 16.",
            "Reporter controls: 2.55 -> 2.78 over the same window.",
            "Interpret as punctate-to-diffuse reporter redistribution, not rupture proof.",
        ],
        1.05,
        6.3,
        5.35,
        1.35,
        12,
    )


def plate_variability_slide(prs):
    slide = blank(prs)
    style_slide(slide, 5)
    add_title(slide, "Plate-level heterogeneity", "All mCherry-valid rows are visible in one heatmap.")
    add_image(slide, ALL / "figures/all_wells_mcherry_ratio_slope_heatmap.png", 0.78, 2.15, 6.0, 3.55)
    add_text(slide, "This view helps choose wells/timepoints for manual ROI review, segmentation validation, and follow-up assays.", 1.05, 6.05, 5.4, 0.55, 11, BLACK)
    add_image(slide, ALL / "figures/all_wells_registration_qc_pass_fraction.png", 1.2, 6.85, 4.9, 1.8)


def single_neuron_slide(prs):
    slide = blank(prs)
    style_slide(slide, 6)
    add_title(slide, "Single-neuron candidate examples", "Local ROI crops make the alignment concept visible.")
    add_image(slide, SINGLE / "figures/E05_vs_F05_single_neuron_mcherry_montage.png", 0.88, 2.28, 5.95, 3.25)
    add_bullets(
        slide,
        [
            "Automatic 488-positive local ROI selection across Days 8/12/16.",
            "Crop-level local registration is applied before mCherry quantification.",
            "Manual same-neuron identity review is still required.",
        ],
        1.05,
        5.95,
        5.5,
        1.25,
        11.5,
    )
    add_image(slide, SINGLE / "figures/E05_vs_F05_single_neuron_mcherry.gif", 1.0, 7.55, 2.55, 1.25)
    add_image(slide, SINGLE / "figures/single_neuron_mcherry_metric_over_time.png", 3.7, 7.25, 2.55, 1.55)
    add_text(slide, "Embedded GIF: E05 vs F05 local ROI mCherry", 1.02, 8.93, 5.1, 0.25, 8, MUTED)


def omezarr_slide(prs):
    slide = blank(prs)
    style_slide(slide, 7)
    add_title(slide, "Why OME-Zarr next", "Metrics are compact; images need a scalable store.")
    add_bullets(
        slide,
        [
            "The all-well run writes compact QC/measurement CSVs rather than duplicating hundreds of image stacks.",
            "OME-Zarr supports chunked, lazy image access for napari and web viewers.",
            "Registered single-neuron ROI stacks were exported as proof-of-format.",
            "Next: plate-level OME-Zarr layout for selected registered wells/days/channels.",
        ],
        1.08,
        2.75,
        5.3,
        2.3,
        12,
    )
    add_text(slide, "Verified OME-Zarr example", 1.12, 5.65, 3.0, 0.35, 15, DARK, bold=True)
    add_bullets(
        slide,
        [
            "Shape: 3 timepoints x 3 channels x 160 x 160 px",
            "Chunks: 1 x 1 x 160 x 160",
            "Location: single_neuron_examples/ome_zarr/",
        ],
        1.15,
        6.15,
        5.1,
        1.1,
        11.5,
    )


def next_steps_slide(prs):
    slide = blank(prs)
    style_slide(slide, 8)
    add_title(slide, "What is needed before a manuscript claim", "This is a strong screening workflow, not final rupture proof.")
    add_bullets(
        slide,
        [
            "Manual ROI review and segmentation/tracking validation for same-neuron claims.",
            "Extend from Days 8/12/16 to all available days: 8, 12, 16, 20, 25, 29, 32, 36, 39.",
            "Orthogonal rupture markers: Galectin-3/8, LAMP1/2 morphology, LysoTracker loss, p62/LC3, LLOMe positive control.",
            "Export selected registered wells to OME-Zarr for scalable review.",
            "Then test condition x time with well/site/neuron-level replicates.",
        ],
        1.08,
        2.55,
        5.35,
        3.4,
        12.5,
    )
    add_text(slide, "Bottom line", 1.1, 6.85, 1.8, 0.35, 16, DARK, bold=True)
    add_text(slide, "The pipeline now turns longitudinal microscopy into all-well quantitative trajectories that can prioritize TMEM106B conditions for deeper lysosome validation.", 1.12, 7.35, 5.3, 0.95, 13, BLACK)
    add_text(slide, "Thank you", 1.12, 8.75, 2.6, 0.55, 26, DARK, bold=True)


def write_speaker_notes() -> None:
    notes = OUT_DIR / "tmem106b_neuron_alignment_speaker_notes.txt"
    notes.write_text(
        """60-second summary:
We built a longitudinal neuron alignment workflow for the TMEM106B imaging dataset. The first E05/F05 run was just a proof-of-pipeline; the current all-well run processed 192 wells across Days 8, 12, and 16, reading 576 ND2 files. Registration used the stable 488 channel so the mCherry phenotype was not used to align itself, and all 576 registration observations passed QC. mCherry puncta/diffuse measurements were computed only for reporter-containing rows E/F/I/J/M/N. The PLD3+TMEM106B+mCherry wells show a stronger increase in diffuse-to-punctate mCherry ratio than reporter controls, consistent with a rupture-like screening phenotype, but not proof of lysosomal rupture.

3-minute summary:
The biological goal is to connect longitudinal light microscopy to the TMEM106B lysosomal fibril/rupture model. The pipeline loads raw ND2 data read-only, registers each well over time using the 488 channel, applies the same shift to 561/mCherry, and measures punctate versus diffuse mCherry signal. We scaled from the initial E05/F05 pilot to all wells in C05:N20 for Days 8, 12, and 16. The all-well run processed 192 wells and 576 files, with 576/576 registration QC pass. Importantly, non-mCherry rows were not treated as zero-puncta samples; they were included only for registration and plate coverage. For the 96 reporter-containing wells, the PLD3+TMEM106B+mCherry condition increased from a mean diffuse/punctate ratio of about 3.26 on Day 8 to 4.53 on Day 16, while PLD3+mCherry reporter controls rose more modestly from 2.55 to 2.78. This is a useful same-well screening metric for punctate-to-diffuse reporter redistribution. Next steps are manual ROI validation, segmentation/tracking, all-day analysis, OME-Zarr plate browsing, and orthogonal rupture assays such as Galectin-3/8, LAMP1/2, LysoTracker, p62/LC3, and LLOMe controls.
""",
        encoding="utf-8",
    )


if __name__ == "__main__":
    main()
